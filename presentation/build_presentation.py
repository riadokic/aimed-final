#!/usr/bin/env python3
"""
AiMED Company Presentation Builder v2
Synced with aimed_pitchdeck.md (copy) + canva_fix.md (design fixes)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ═══════════════════════════════════════════════════════════
# BRAND CONSTANTS (synced with canva_fix.md)
# ═══════════════════════════════════════════════════════════

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
NEAR_BLACK = RGBColor(0x0D, 0x1B, 0x2A)
ZINC_700 = RGBColor(0x3F, 0x3F, 0x46)  # Body text color (fix: was zinc-500)
ZINC_600 = RGBColor(0x52, 0x52, 0x5B)  # FAQ answer text
ZINC_500 = RGBColor(0x71, 0x71, 0x7A)  # Labels (fix: was zinc-400)
ZINC_400 = RGBColor(0xA1, 0xA1, 0xAA)
ZINC_300 = RGBColor(0xD4, 0xD4, 0xD8)  # Faded H1 second line
ZINC_200 = RGBColor(0xE4, 0xE4, 0xE7)
ZINC_100 = RGBColor(0xF4, 0xF4, 0xF5)  # Card borders
ZINC_50 = RGBColor(0xFA, 0xFA, 0xFA)   # Card backgrounds
BLUE_600 = RGBColor(0x25, 0x63, 0xEB)
BLUE_50 = RGBColor(0xEF, 0xF6, 0xFF)
DARK_SURFACE = RGBColor(0x11, 0x11, 0x11)

# Slide dimensions (16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Margins — canva_fix: min 48px content-to-edge
MARGIN = Inches(0.6)

FONT_FAMILY = "Inter"

# Fix constants from canva_fix.md
H1_SIZE_LIGHT = Pt(44)       # Fix: was 64/56/48 → all light slides 44pt
H1_LINE_HEIGHT = 1.05        # Fix: was 0.92
H1_SIZE_DARK = Pt(72)        # Title/CTA slides keep 72pt
BODY_SIZE = Pt(15)            # Fix: was 16-17pt → 15-16pt
BODY_COLOR = ZINC_700         # Fix: was ZINC_500
BODY_LINE_HEIGHT = 1.65       # Fix: was 1.6-1.7
LABEL_SIZE = Pt(10)           # Fix: was 11pt
LABEL_COLOR = ZINC_500        # Fix: was ZINC_400
LABEL_SPACING = 180           # Fix: was 200
LOGO_WORDMARK_SIZE = Pt(14)   # Fix: was 20pt
LOGO_SQUARE_SIZE = Inches(0.29)  # ~28px
LOGO_AI_SIZE = Pt(12)

# ═══════════════════════════════════════════════════════════
# HELPER FUNCTIONS
# ═══════════════════════════════════════════════════════════

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rounded_rect(slide, left, top, width, height, fill_color=None,
                     border_color=None, border_width=Pt(1)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.adjustments[0] = 0.05

    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()

    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()

    return shape


def add_text_box(slide, left, top, width, height, text="", font_size=Pt(16),
                 font_color=NEAR_BLACK, bold=False, alignment=PP_ALIGN.LEFT,
                 line_spacing=1.6, anchor=MSO_ANCHOR.TOP, all_caps=False,
                 spacing=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except:
        pass

    p = tf.paragraphs[0]
    p.text = text.upper() if all_caps else text
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing:
        p.line_spacing = line_spacing

    run = p.runs[0] if p.runs else p.add_run()
    if not p.runs:
        run.text = text.upper() if all_caps else text
    run.font.size = font_size
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.name = FONT_FAMILY

    if spacing:
        rPr = run._r.find('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
        if rPr is not None:
            rPr.set('spc', str(spacing))

    return txBox


def add_multiline_text(slide, left, top, width, height, lines, alignment=PP_ALIGN.LEFT,
                       anchor=MSO_ANCHOR.TOP, line_spacing=1.6):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except:
        pass

    for i, line_def in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        text = line_def.get('text', '')
        is_caps = line_def.get('all_caps', False)
        p.text = ""
        p.alignment = alignment
        p.space_after = line_def.get('space_after', Pt(0))
        p.space_before = line_def.get('space_before', Pt(0))
        p.line_spacing = line_def.get('line_spacing', line_spacing)

        run = p.add_run()
        run.text = text.upper() if is_caps else text
        run.font.size = line_def.get('font_size', Pt(16))
        run.font.color.rgb = line_def.get('font_color', NEAR_BLACK)
        run.font.bold = line_def.get('bold', False)
        run.font.name = FONT_FAMILY

        spacing = line_def.get('spacing', None)
        if spacing:
            rPr = run._r.find('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
            if rPr is not None:
                rPr.set('spc', str(spacing))

    return txBox


def add_page_number(slide, number, total, is_dark=False):
    add_text_box(
        slide,
        left=SLIDE_W - Inches(1.5),
        top=SLIDE_H - Inches(0.5),
        width=Inches(1.2),
        height=Inches(0.3),
        text=f"{number:02d} / {total:02d}",
        font_size=Pt(11),
        font_color=ZINC_400,
        alignment=PP_ALIGN.RIGHT
    )


def add_logo(slide, is_dark=False):
    """Add AiMED logo — fix: 28x28 square, 14pt wordmark, consistent position."""
    logo_color = WHITE if is_dark else NEAR_BLACK
    sq = add_rounded_rect(
        slide, left=MARGIN, top=MARGIN,
        width=LOGO_SQUARE_SIZE, height=LOGO_SQUARE_SIZE,
        fill_color=WHITE if is_dark else BLACK
    )
    # "Ai" monogram
    add_text_box(
        slide,
        left=MARGIN + Emu(Inches(0.01)),
        top=MARGIN - Emu(Inches(0.01)),
        width=LOGO_SQUARE_SIZE,
        height=LOGO_SQUARE_SIZE,
        text="Ai",
        font_size=LOGO_AI_SIZE,
        font_color=BLACK if is_dark else WHITE,
        bold=True,
        alignment=PP_ALIGN.CENTER,
        line_spacing=1.0
    )

    # Wordmark — fix: 14pt (was 20pt)
    add_text_box(
        slide,
        left=MARGIN + LOGO_SQUARE_SIZE + Inches(0.1),
        top=MARGIN - Inches(0.01),
        width=Inches(1.2),
        height=Inches(0.35),
        text="AiMED",
        font_size=LOGO_WORDMARK_SIZE,
        font_color=logo_color,
        bold=True,
        alignment=PP_ALIGN.LEFT,
        line_spacing=1.0
    )


def add_label(slide, text, left=None, top=None):
    """Fix: 10pt, Bold, zinc-500, spacing +180, gap 10px before H1."""
    if left is None:
        left = MARGIN
    if top is None:
        top = MARGIN + Inches(0.7)

    return add_text_box(
        slide, left=left, top=top,
        width=Inches(8), height=Inches(0.25),
        text=text,
        font_size=LABEL_SIZE,
        font_color=LABEL_COLOR,
        bold=True,
        alignment=PP_ALIGN.LEFT,
        all_caps=True,
        spacing=LABEL_SPACING,
        line_spacing=1.0
    )


def add_badge_pill(slide, left, top, text, border_color=WHITE, text_color=WHITE,
                   fill_color=None, width=Inches(2.0), height=Inches(0.35)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.adjustments[0] = 0.5

    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()

    shape.line.color.rgb = border_color
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    try:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except:
        pass
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text.upper()
    run.font.size = Pt(10)
    run.font.color.rgb = text_color
    run.font.bold = True
    run.font.name = FONT_FAMILY
    return shape


# ═══════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════

TOTAL_SLIDES = 10


def build_slide_01(prs):
    """SLIDE 01 — NASLOVNI SLAJD (Dark) — no changes needed per fix doc."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BLACK)
    add_logo(slide, is_dark=True)

    # Company label
    add_text_box(
        slide, left=MARGIN, top=MARGIN + Inches(0.4),
        width=Inches(3), height=Inches(0.25),
        text="CEE-MED D.O.O.",
        font_size=Pt(11), font_color=ZINC_400,
        bold=True, all_caps=True, spacing=200, line_spacing=1.0
    )

    # Small label
    add_text_box(
        slide, left=Inches(0), top=Inches(2.2),
        width=SLIDE_W, height=Inches(0.3),
        text="INTELIGENTNA MEDICINSKA DOKUMENTACIJA",
        font_size=Pt(12), font_color=ZINC_500,
        bold=False, alignment=PP_ALIGN.CENTER,
        all_caps=True, spacing=200, line_spacing=1.0
    )

    # H1 — dark slides keep 72pt
    add_multiline_text(
        slide, left=Inches(0), top=Inches(2.6),
        width=SLIDE_W, height=Inches(2.2),
        alignment=PP_ALIGN.CENTER, line_spacing=0.92,
        lines=[
            {'text': 'Vi diktirate,', 'font_size': Pt(72), 'font_color': WHITE, 'bold': True, 'line_spacing': 0.92},
            {'text': 'AiMED piše.', 'font_size': Pt(72), 'font_color': WHITE, 'bold': True, 'line_spacing': 0.92},
        ]
    )

    # Subtitle
    add_multiline_text(
        slide, left=Inches(2), top=Inches(4.6),
        width=Inches(9.333), height=Inches(1.0),
        alignment=PP_ALIGN.CENTER,
        lines=[
            {'text': 'Inteligentni AI asistent koji pretvara glasovni diktat ljekara', 'font_size': Pt(20), 'font_color': ZINC_400, 'line_spacing': 1.5},
            {'text': 'u strukturirani medicinski nalaz \u2014 za manje od 60 sekundi.', 'font_size': Pt(20), 'font_color': ZINC_400, 'line_spacing': 1.5},
        ]
    )

    # Badges
    badge_y = Inches(6.0)
    badge_w = Inches(2.4)
    badge_gap = Inches(0.3)
    total_badges_w = badge_w * 3 + badge_gap * 2
    start_x = (SLIDE_W - total_badges_w) / 2

    for i, badge_text in enumerate(["ISO 27001:2022", "GDPR USKLAĐENO", "ZZLP BiH"]):
        add_badge_pill(
            slide,
            left=int(start_x) + int((badge_w + badge_gap) * i),
            top=badge_y, text=badge_text,
            border_color=ZINC_700, text_color=ZINC_400,
            width=badge_w, height=Inches(0.38)
        )

    add_page_number(slide, 1, TOTAL_SLIDES, is_dark=True)


def build_slide_02(prs):
    """SLIDE 02 — PROBLEM (Light)
    Fix: H1 44pt, body zinc-700 15pt, copy from pitchdeck.md
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_logo(slide, is_dark=False)
    add_label(slide, "Problem koji prepoznajete")

    # H1 — fix: 44pt (was 64pt), expanded text box
    add_multiline_text(
        slide, left=MARGIN, top=Inches(1.1),
        width=Inches(5.8), height=Inches(1.6),
        line_spacing=H1_LINE_HEIGHT,
        lines=[
            {'text': 'Papiri umjesto', 'font_size': H1_SIZE_LIGHT, 'font_color': NEAR_BLACK, 'bold': True, 'line_spacing': H1_LINE_HEIGHT},
            {'text': 'pacijenata.', 'font_size': H1_SIZE_LIGHT, 'font_color': ZINC_300, 'bold': True, 'line_spacing': H1_LINE_HEIGHT},
        ]
    )

    # Body — synced from aimed_pitchdeck.md, fix: zinc-700 15pt
    add_multiline_text(
        slide, left=MARGIN, top=Inches(2.8),
        width=Inches(5.5), height=Inches(3.0),
        lines=[
            {'text': 'Za svakog pacijenta postoji procedura koja se ponavlja: otvori program, upiši anamnezu, provjeri dijagnozu po MKB-10 listi, odaberi lijek iz Registra, formatiraj nalaz, printaj i potpiši.', 'font_size': BODY_SIZE, 'font_color': BODY_COLOR, 'line_spacing': BODY_LINE_HEIGHT},
            {'text': '', 'font_size': Pt(8), 'font_color': BODY_COLOR, 'line_spacing': 1.0},
            {'text': 'Svaki put. Za svaku osobu.', 'font_size': BODY_SIZE, 'font_color': BODY_COLOR, 'bold': True, 'line_spacing': BODY_LINE_HEIGHT},
            {'text': '', 'font_size': Pt(8), 'font_color': BODY_COLOR, 'line_spacing': 1.0},
            {'text': 'Zbrojeno: to su stotine sati godišnje koje ne provodite sa pacijentima.', 'font_size': BODY_SIZE, 'font_color': BODY_COLOR, 'line_spacing': BODY_LINE_HEIGHT},
        ]
    )

    # Right column — 3 stat cards (fix: zinc-50 bg, zinc-100 border, 24px padding)
    card_left = Inches(7.0)
    card_width = Inches(5.7)
    card_height = Inches(1.55)
    card_gap = Inches(0.2)

    # Card 1
    card1_top = Inches(1.0)
    add_rounded_rect(slide, card_left, card1_top, card_width, card_height,
                     fill_color=ZINC_50, border_color=ZINC_100)
    add_multiline_text(
        slide, card_left + Inches(0.35), card1_top + Inches(0.25),
        card_width - Inches(0.7), card_height - Inches(0.4),
        lines=[
            {'text': '~27%', 'font_size': Pt(48), 'font_color': NEAR_BLACK, 'bold': True, 'line_spacing': 1.1},
            {'text': 'radnog vremena na dokumentaciju (~15 h sedmično) ¹', 'font_size': Pt(13), 'font_color': ZINC_600, 'line_spacing': 1.4, 'space_before': Pt(4)},
        ]
    )

    # Card 2 — BLACK accent (fix: no border, no shadow, 28px padding)
    card2_top = card1_top + card_height + card_gap
    add_rounded_rect(slide, card_left, card2_top, card_width, card_height,
                     fill_color=BLACK)
    add_multiline_text(
        slide, card_left + Inches(0.4), card2_top + Inches(0.25),
        card_width - Inches(0.8), card_height - Inches(0.4),
        lines=[
            {'text': '+1,8 h/dan', 'font_size': Pt(48), 'font_color': WHITE, 'bold': True, 'line_spacing': 1.1},
            {'text': 'dokumentacije izvan ordinirajućeg vremena ²', 'font_size': Pt(13), 'font_color': ZINC_400, 'line_spacing': 1.4, 'space_before': Pt(4)},
        ]
    )

    # Card 3
    card3_top = card2_top + card_height + card_gap
    add_rounded_rect(slide, card_left, card3_top, card_width, card_height,
                     fill_color=ZINC_50, border_color=ZINC_100)
    add_multiline_text(
        slide, card_left + Inches(0.35), card3_top + Inches(0.25),
        card_width - Inches(0.7), card_height - Inches(0.4),
        lines=[
            {'text': '53%', 'font_size': Pt(48), 'font_color': NEAR_BLACK, 'bold': True, 'line_spacing': 1.1},
            {'text': 'ljekara s burnout simptomima \u2014 administracija faktor #1 (Medscape 2023) ³', 'font_size': Pt(13), 'font_color': ZINC_600, 'line_spacing': 1.4, 'space_before': Pt(4)},
        ]
    )

    add_page_number(slide, 2, TOTAL_SLIDES)


def build_slide_03(prs):
    """SLIDE 03 — KONTEKST BiH (Light)
    Fix: H1 44pt, body zinc-700, copy from pitchdeck.md
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_logo(slide, is_dark=False)
    add_label(slide, "Realnost zdravstvenog sistema u BiH")

    # H1 — fix: 44pt
    add_multiline_text(
        slide, MARGIN, Inches(1.1), Inches(8), Inches(1.6),
        line_spacing=H1_LINE_HEIGHT,
        lines=[
            {'text': 'Premalo ljekara.', 'font_size': H1_SIZE_LIGHT, 'font_color': NEAR_BLACK, 'bold': True, 'line_spacing': H1_LINE_HEIGHT},
            {'text': 'Previše papira.', 'font_size': H1_SIZE_LIGHT, 'font_color': ZINC_300, 'bold': True, 'line_spacing': H1_LINE_HEIGHT},
        ]
    )

    # Comparison bars
    bar_top = Inches(2.9)
    bar_left = MARGIN

    # EU bar
    add_text_box(slide, bar_left, bar_top, Inches(2.5), Inches(0.35),
                 text="\U0001F1EA\U0001F1FA  EU prosjek", font_size=Pt(14), font_color=NEAR_BLACK, bold=True)
    eu_bar = add_rounded_rect(slide, bar_left + Inches(2.8), bar_top + Inches(0.02),
                               Inches(6.5), Inches(0.3), fill_color=NEAR_BLACK)
    eu_bar.adjustments[0] = 0.5
    add_text_box(slide, bar_left + Inches(9.5), bar_top, Inches(3.5), Inches(0.35),
                 text="3,9 ljekara / 1.000 stan. (Eurostat 2023)", font_size=Pt(13), font_color=ZINC_600)

    # BiH bar
    bih_top = bar_top + Inches(0.6)
    add_text_box(slide, bar_left, bih_top, Inches(2.5), Inches(0.35),
                 text="\U0001F1E7\U0001F1E6  BiH", font_size=Pt(14), font_color=NEAR_BLACK, bold=True)
    bih_bar = add_rounded_rect(slide, bar_left + Inches(2.8), bih_top + Inches(0.02),
                                Inches(3.5), Inches(0.3), fill_color=ZINC_300)
    bih_bar.adjustments[0] = 0.5
    add_text_box(slide, bar_left + Inches(6.5), bih_top, Inches(5.5), Inches(0.35),
                 text="~2,1 ljekara / 1.000 stan.  \u2190 skoro 50% manje", font_size=Pt(13), font_color=ZINC_600)

    # Bottom stats + narrative (synced from pitchdeck)
    col_top = Inches(4.4)

    # Stat 1
    add_multiline_text(
        slide, MARGIN, col_top, Inches(3.5), Inches(1.0),
        lines=[
            {'text': '~5.800', 'font_size': Pt(48), 'font_color': NEAR_BLACK, 'bold': True, 'line_spacing': 1.1},
            {'text': 'aktivnih ljekara u BiH', 'font_size': Pt(13), 'font_color': ZINC_600, 'line_spacing': 1.4},
        ]
    )

    # Stat 2
    add_multiline_text(
        slide, Inches(4.2), col_top, Inches(3.5), Inches(1.0),
        lines=[
            {'text': '~750', 'font_size': Pt(48), 'font_color': NEAR_BLACK, 'bold': True, 'line_spacing': 1.1},
            {'text': 'ljekara porod. medicine \u2014 deficit ~50% (ZZJZFBIH / RTVBN 2023)', 'font_size': Pt(13), 'font_color': ZINC_600, 'line_spacing': 1.4},
        ]
    )

    # BLACK accent card — synced narrative from pitchdeck
    accent_left = Inches(8.2)
    accent_w = Inches(4.5)
    accent_h = Inches(2.8)
    add_rounded_rect(slide, accent_left, col_top, accent_w, accent_h, fill_color=BLACK)
    add_multiline_text(
        slide, accent_left + Inches(0.4), col_top + Inches(0.35),
        accent_w - Inches(0.8), accent_h - Inches(0.7),
        alignment=PP_ALIGN.LEFT,
        lines=[
            {'text': 'Svaki ljekar u BiH de facto radi za dvojicu.', 'font_size': Pt(16), 'font_color': WHITE, 'bold': True, 'line_spacing': 1.6},
            {'text': '', 'font_size': Pt(6), 'font_color': WHITE, 'line_spacing': 1.0},
            {'text': 'A zatim provede trećinu tog pretrpanog radnog dana ne s pacijentima \u2014 nego na administraciji.', 'font_size': Pt(15), 'font_color': ZINC_400, 'line_spacing': 1.6},
            {'text': '', 'font_size': Pt(8), 'font_color': WHITE, 'line_spacing': 1.0},
            {'text': 'Kad svaki sat broji \u2014 svaki minut koji AI preuzme direktno se pretvara u vrijeme s pacijentom.', 'font_size': Pt(15), 'font_color': WHITE, 'bold': True, 'line_spacing': 1.6},
        ]
    )

    add_page_number(slide, 3, TOTAL_SLIDES)


def build_slide_04(prs):
    """SLIDE 04 — RJEŠENJE (Light)
    Fix: H1 44pt, body from pitchdeck (web-app explanation), pills from pitchdeck
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_logo(slide, is_dark=False)
    add_label(slide, "Naše rješenje")

    # H1 — fix: 44pt
    add_multiline_text(
        slide, MARGIN, Inches(1.1), Inches(7.5), Inches(1.6),
        line_spacing=H1_LINE_HEIGHT,
        lines=[
            {'text': 'Govorite slobodno.', 'font_size': H1_SIZE_LIGHT, 'font_color': NEAR_BLACK, 'bold': True, 'line_spacing': H1_LINE_HEIGHT},
            {'text': 'AiMED strukturira sve ostalo.', 'font_size': H1_SIZE_LIGHT, 'font_color': ZINC_300, 'bold': True, 'line_spacing': H1_LINE_HEIGHT},
        ]
    )

    # Body — synced from pitchdeck (warmer, colleague tone)
    add_multiline_text(
        slide, MARGIN, Inches(2.8), Inches(5.8), Inches(2.8),
        lines=[
            {'text': 'AiMED je web-aplikacija kojoj pristupate iz vašeg preglednika \u2014 iste kao što otvarate email ili neku web stranicu.', 'font_size': BODY_SIZE, 'font_color': BODY_COLOR, 'line_spacing': BODY_LINE_HEIGHT},
            {'text': '', 'font_size': Pt(6), 'font_color': BODY_COLOR, 'line_spacing': 1.0},
            {'text': 'Ne treba instalacija, ne treba poseban kompjuter.', 'font_size': BODY_SIZE, 'font_color': BODY_COLOR, 'line_spacing': BODY_LINE_HEIGHT},
            {'text': '', 'font_size': Pt(6), 'font_color': BODY_COLOR, 'line_spacing': 1.0},
            {'text': 'Govorite na bosanskom, hrvatskom ili srpskom, potpuno prirodno, kao što biste diktirali sestri ili kolegi.', 'font_size': BODY_SIZE, 'font_color': BODY_COLOR, 'line_spacing': BODY_LINE_HEIGHT},
            {'text': '', 'font_size': Pt(6), 'font_color': BODY_COLOR, 'line_spacing': 1.0},
            {'text': 'AiMED to sluša, prepoznaje medicinske pojmove i za manje od jedne minute složi kompletan nalaz spreman za potpis.', 'font_size': BODY_SIZE, 'font_color': BODY_COLOR, 'line_spacing': BODY_LINE_HEIGHT},
        ]
    )

    # 3 Feature pills — synced from pitchdeck
    pill_top = Inches(5.7)
    pill_h = Inches(1.2)
    pill_w = Inches(3.8)
    pill_gap = Inches(0.25)

    pills = [
        ("\U0001F399  BHS JEZIK", "Jedini ASR engine finetunovan na lokalni medicinski rječnik. Razumije \"Sumamed\", \"KCUS\", lokalne termine."),
        ("\u26A1  BRZINA & PRECIZNOST", "98,6% preciznost (EUSIPCO 2022). Diktat 2 min \u2192 nalaz < 60 sek. Ušteda 25\u201340%."),
        ("\U0001F512  PRIVATNOST PO DIZAJNU", "Glasovni zapisi brišu se odmah po obradi. Nema pohrane audio snimaka."),
    ]

    for i, (title, desc) in enumerate(pills):
        px = MARGIN + int((pill_w + pill_gap) * i)
        add_rounded_rect(slide, px, pill_top, pill_w, pill_h,
                         fill_color=ZINC_50, border_color=ZINC_100)
        add_multiline_text(
            slide, px + Inches(0.25), pill_top + Inches(0.15),
            pill_w - Inches(0.5), pill_h - Inches(0.25),
            lines=[
                {'text': title, 'font_size': Pt(11), 'font_color': NEAR_BLACK, 'bold': True, 'line_spacing': 1.2, 'all_caps': True, 'spacing': 150},
                {'text': desc, 'font_size': Pt(12), 'font_color': ZINC_600, 'line_spacing': 1.5, 'space_before': Pt(6)},
            ]
        )

    # Right visual — flow diagram
    flow_left = Inches(8.0)
    flow_top = Inches(1.3)
    box_w = Inches(4.7)
    box_h = Inches(0.9)
    box_gap = Inches(0.25)

    boxes = [
        ("\U0001F399 Glas ljekara", "2 min diktat", ZINC_50, NEAR_BLACK),
        ("AiMED AI", "", BLACK, WHITE),
        ("\U0001F4C4 Potpun nalaz", "PDF + Word + Copy", ZINC_50, NEAR_BLACK),
    ]

    for i, (title, sub, bg, fg) in enumerate(boxes):
        by = flow_top + int((box_h + box_gap) * i)
        add_rounded_rect(slide, flow_left, by, box_w, box_h, fill_color=bg,
                         border_color=ZINC_100 if bg != BLACK else None)
        content_lines = [
            {'text': title, 'font_size': Pt(16), 'font_color': fg, 'bold': True, 'line_spacing': 1.2},
        ]
        if sub:
            content_lines.append({'text': sub, 'font_size': Pt(12), 'font_color': ZINC_400 if bg == BLACK else ZINC_600, 'line_spacing': 1.3})
        add_multiline_text(
            slide, flow_left + Inches(0.3), by + Inches(0.1),
            box_w - Inches(0.6), box_h - Inches(0.2),
            alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
            lines=content_lines
        )
        if i < 2:
            arrow_y = by + box_h + Inches(0.02)
            add_text_box(slide, flow_left, arrow_y, box_w, Inches(0.2),
                         text="\u2193", font_size=Pt(18), font_color=ZINC_300,
                         alignment=PP_ALIGN.CENTER, line_spacing=1.0)

    add_page_number(slide, 4, TOTAL_SLIDES)


def build_slide_05(prs):
    """SLIDE 05 — KAKO RADI (Light)
    Fix: H1 44pt, card bodies from pitchdeck (conversational)
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_logo(slide, is_dark=False)
    add_label(slide, "Workflow \u2014 korak po korak")

    # H1 — fix: 44pt
    add_text_box(slide, MARGIN, Inches(1.1), Inches(8), Inches(0.8),
                 text="\u010cetiri koraka. Jedan klik.", font_size=H1_SIZE_LIGHT,
                 font_color=NEAR_BLACK, bold=True, line_spacing=H1_LINE_HEIGHT)

    # 4 step cards — bodies synced from pitchdeck
    card_top = Inches(2.3)
    card_w = Inches(2.9)
    card_h = Inches(3.6)
    card_gap = Inches(0.2)

    steps = [
        ("01", "Snimanje",
         "Otvorite AiMED u pregledniku i pritisnite dugme za snimanje. Govorite slobodno, kao što razgovarate s pacijentom.",
         False),
        ("02", "Obrada u realnom\nvremenu",
         "AiMED razumije naš medicinski rječnik: dijagnoze, lijekove, latinske nazive. Preciznost: 98,6% (EUSIPCO 2022).",
         False),
        ("03", "AI Strukturiranje",
         "Sistem sam prepoznaje šta je anamneza, šta dijagnoza, šta terapija \u2014 i sve svrstava na pravo mjesto u nalazu.",
         True),
        ("04", "Pregled & Export",
         "Pregledate nalaz, dodate potpis i pečat, i šaljete pacijentu. Ili kopirate u vaš bolnički program jednim klikom.",
         False),
    ]

    for i, (num, title, body, is_accent) in enumerate(steps):
        cx = int(MARGIN) + int((card_w + card_gap) * i)
        bg = BLACK if is_accent else ZINC_50
        border = None if is_accent else ZINC_100

        add_rounded_rect(slide, cx, card_top, card_w, card_h,
                         fill_color=bg, border_color=border)

        # Step number
        num_color = RGBColor(0x27, 0x27, 0x2A) if is_accent else ZINC_200
        add_text_box(slide, cx + Inches(0.25), card_top + Inches(0.2),
                     Inches(1.5), Inches(0.8),
                     text=num, font_size=Pt(48), font_color=num_color,
                     bold=True, line_spacing=1.0)

        # Title
        add_text_box(slide, cx + Inches(0.25), card_top + Inches(1.0),
                     card_w - Inches(0.5), Inches(0.7),
                     text=title, font_size=Pt(16),
                     font_color=WHITE if is_accent else NEAR_BLACK,
                     bold=True, line_spacing=1.2)

        # Body — fix: 14pt Regular zinc-700 (accent: zinc-400)
        add_text_box(slide, cx + Inches(0.25), card_top + Inches(1.8),
                     card_w - Inches(0.5), Inches(1.7),
                     text=body, font_size=Pt(13),
                     font_color=ZINC_400 if is_accent else ZINC_600,
                     line_spacing=1.5)

    # Bottom strip
    strip_top = Inches(6.2)
    strip_h = Inches(0.7)
    add_rounded_rect(slide, MARGIN, strip_top, SLIDE_W - MARGIN * 2, strip_h,
                     fill_color=ZINC_50, border_color=ZINC_100)

    strip_items = [
        ("Ručno: 20\u201330 minuta", ZINC_600),
        ("\u2192", ZINC_300),
        ("S AiMED-om: ~3 minute", NEAR_BLACK),
        ("\u2192", ZINC_300),
        ("Ušteda: 25\u201340% vremena", NEAR_BLACK),
    ]
    item_w = Inches(2.4)
    for i, (text, color) in enumerate(strip_items):
        ix = int(MARGIN + Inches(0.1)) + int(item_w * i)
        add_text_box(slide, ix, strip_top + Inches(0.1),
                     item_w, strip_h - Inches(0.2),
                     text=text, font_size=Pt(14), font_color=color,
                     bold=True, alignment=PP_ALIGN.CENTER, line_spacing=1.0)

    add_page_number(slide, 5, TOTAL_SLIDES)


def build_slide_06(prs):
    """SLIDE 06 — FUNKCIJE (Light)
    Fix: H1 44pt, expanded text box so "vaše vrijeme" is visible,
    card bodies from canva_fix copy rewrites
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_logo(slide, is_dark=False)
    add_label(slide, "Šta dobijate")

    # H1 — fix: 44pt, expanded height so "vaše vrijeme" is fully visible
    add_multiline_text(
        slide, MARGIN, Inches(1.1), Inches(8), Inches(1.6),
        line_spacing=H1_LINE_HEIGHT,
        lines=[
            {'text': 'Dva načina rada.', 'font_size': H1_SIZE_LIGHT, 'font_color': NEAR_BLACK, 'bold': True, 'line_spacing': H1_LINE_HEIGHT},
            {'text': 'Jedan cilj: vaše vrijeme.', 'font_size': H1_SIZE_LIGHT, 'font_color': ZINC_300, 'bold': True, 'line_spacing': H1_LINE_HEIGHT},
        ]
    )

    # Bento grid
    grid_top = Inches(2.9)
    grid_left = MARGIN
    big_w = Inches(7.8)
    small_w = Inches(4.2)
    row_h = Inches(2.0)
    gap = Inches(0.2)

    # TOP LEFT — Novi snimak (fix: 14pt zinc-600 body)
    add_rounded_rect(slide, grid_left, grid_top, big_w, row_h,
                     fill_color=ZINC_50, border_color=ZINC_100)
    add_badge_pill(slide, grid_left + Inches(0.3), grid_top + Inches(0.2),
                   "NAJČEŠĆE KORIŠTENO", border_color=BLUE_600, text_color=BLUE_600,
                   fill_color=BLUE_50, width=Inches(2.5), height=Inches(0.28))
    add_multiline_text(
        slide, grid_left + Inches(0.3), grid_top + Inches(0.6),
        big_w - Inches(0.6), row_h - Inches(0.7),
        lines=[
            {'text': '\U0001F399  Novi snimak', 'font_size': Pt(22), 'font_color': NEAR_BLACK, 'bold': True, 'line_spacing': 1.3},
            {'text': "Kliknete 'Snimi' i počnete govoriti. Ne morate ništa unaprijed planirati \u2014 sistem strukturira sve iz vašeg glasa, automatski.", 'font_size': Pt(14), 'font_color': ZINC_600, 'line_spacing': 1.5, 'space_before': Pt(4)},
        ]
    )

    # TOP RIGHT — Ažuriranje nalaza (fix: 13pt zinc-600)
    sr_left = grid_left + big_w + gap
    add_rounded_rect(slide, sr_left, grid_top, small_w, row_h,
                     fill_color=ZINC_50, border_color=ZINC_100)
    add_multiline_text(
        slide, sr_left + Inches(0.3), grid_top + Inches(0.25),
        small_w - Inches(0.6), row_h - Inches(0.5),
        lines=[
            {'text': '\U0001F504  Ažuriranje nalaza', 'font_size': Pt(18), 'font_color': NEAR_BLACK, 'bold': True, 'line_spacing': 1.3},
            {'text': 'Imate pacijenta na kontroli? Recite samo što se promijenilo. AiMED ažurira samo te dijelove, ostalo ostaje netaknuto.', 'font_size': Pt(13), 'font_color': ZINC_600, 'line_spacing': 1.5, 'space_before': Pt(4)},
            {'text': '', 'font_size': Pt(4), 'font_color': ZINC_600, 'line_spacing': 1.0},
            {'text': '"Zamijeni Brufen sa Sumamedom 500mg, kontrola za 10 dana."', 'font_size': Pt(12), 'font_color': ZINC_500, 'line_spacing': 1.4},
        ]
    )

    # BOTTOM LEFT — Privatnost (BLACK) (fix: 15pt white, zinc-400 body)
    row2_top = grid_top + row_h + gap
    add_rounded_rect(slide, grid_left, row2_top, big_w, row_h, fill_color=BLACK)
    add_multiline_text(
        slide, grid_left + Inches(0.4), row2_top + Inches(0.35),
        big_w - Inches(0.8), row_h - Inches(0.6),
        lines=[
            {'text': '\U0001F6E1  Privatnost po dizajnu.', 'font_size': Pt(22), 'font_color': WHITE, 'bold': True, 'line_spacing': 1.3},
            {'text': 'Sve što izgovorite se obradi i odmah briše s naših servera. Nijedan audio zapis ne ostaje pohranjen. Nikome. GDPR usklađenost Član 9.', 'font_size': Pt(15), 'font_color': ZINC_400, 'line_spacing': 1.5, 'space_before': Pt(6)},
        ]
    )

    # BOTTOM RIGHT — Šabloni (fix: 13pt zinc-600)
    add_rounded_rect(slide, sr_left, row2_top, small_w, row_h,
                     fill_color=ZINC_50, border_color=ZINC_100)
    add_multiline_text(
        slide, sr_left + Inches(0.3), row2_top + Inches(0.3),
        small_w - Inches(0.6), row_h - Inches(0.5),
        lines=[
            {'text': '\U0001F4C4  Rad sa šablonima', 'font_size': Pt(18), 'font_color': NEAR_BLACK, 'bold': True, 'line_spacing': 1.3},
            {'text': 'Imate obrazac koji koristite godinama? Učitate ga jednom. AiMED ga popunjava iz vašeg diktata, s vašim brandingom.', 'font_size': Pt(13), 'font_color': ZINC_600, 'line_spacing': 1.5, 'space_before': Pt(4)},
        ]
    )

    add_page_number(slide, 6, TOTAL_SLIDES)


def build_slide_07(prs):
    """SLIDE 07 — SIGURNOST (Light)
    Fix: H1 40pt (more content), 24px gap to list, Q&A card synced from pitchdeck,
    added RBAC item from pitchdeck
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_logo(slide, is_dark=False)
    add_label(slide, "Sigurnost & Usklađenost")

    # H1 — fix: 40pt (this slide has more content)
    add_multiline_text(
        slide, MARGIN, Inches(1.1), Inches(7), Inches(1.4),
        line_spacing=H1_LINE_HEIGHT,
        lines=[
            {'text': 'Privatnost po dizajnu.', 'font_size': Pt(40), 'font_color': NEAR_BLACK, 'bold': True, 'line_spacing': H1_LINE_HEIGHT},
            {'text': 'Ne kao afterthought.', 'font_size': Pt(40), 'font_color': ZINC_300, 'bold': True, 'line_spacing': H1_LINE_HEIGHT},
        ]
    )

    # Security list — fix: 14pt Regular zinc-700, 12px between items
    items_top = Inches(2.8)
    items_left = MARGIN
    items_width = Inches(6.5)
    item_h = Inches(0.58)

    # Synced items from pitchdeck (added RBAC)
    security_items = [
        ("\U0001F512  GDPR", "Uredba EU 2016/679, Član 9 (posebna kategorija: medicinski podaci)"),
        ("\U0001F512  ZZLP BiH", "Zakon o zaštiti ličnih podataka BiH (saglasnost za posebne kategorije)"),
        ("\U0001F512  ISO 27001:2022", "Certifikacija u tijeku"),
        ("\U0001F512  Enkripcija", "AES-256 at rest \u00b7 TLS 1.3 in transit"),
        ("\U0001F512  Nema audio pohrane", "Zapisi se brišu odmah po obradi"),
        ("\U0001F512  RBAC pristup", "Svaki ljekar vidi samo svoje nalaze"),
        ("\U0001F512  EU infrastruktura", "Svi serveri unutar EU granica"),
    ]

    for i, (title, desc) in enumerate(security_items):
        iy = items_top + int(item_h * i)
        if i > 0:
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, items_left, iy - Inches(0.02), items_width, Pt(1)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = ZINC_100
            line.line.fill.background()

        add_multiline_text(
            slide, items_left, iy + Inches(0.04), items_width, item_h - Inches(0.08),
            lines=[
                {'text': title, 'font_size': Pt(14), 'font_color': NEAR_BLACK, 'bold': True, 'line_spacing': 1.2},
                {'text': desc, 'font_size': Pt(12), 'font_color': ZINC_700, 'line_spacing': 1.3},
            ]
        )

    # Right BLACK Q&A card — synced Q/A from pitchdeck + canva_fix
    qa_left = Inches(7.8)
    qa_top = Inches(2.6)
    qa_w = Inches(4.9)
    qa_h = Inches(4.2)
    add_rounded_rect(slide, qa_left, qa_top, qa_w, qa_h, fill_color=BLACK)

    add_multiline_text(
        slide, qa_left + Inches(0.4), qa_top + Inches(0.4),
        qa_w - Inches(0.8), qa_h - Inches(0.8),
        lines=[
            {'text': 'NAJČEŠĆE PITANJE:', 'font_size': Pt(10), 'font_color': ZINC_400, 'bold': True, 'line_spacing': 1.2, 'all_caps': True, 'spacing': LABEL_SPACING},
            {'text': '', 'font_size': Pt(10), 'font_color': ZINC_400, 'line_spacing': 1.0},
            {'text': '"Hoće li AiMED čuvati ili dijeliti podatke mojih pacijenata?"', 'font_size': Pt(18), 'font_color': WHITE, 'line_spacing': 1.5, 'space_before': Pt(8)},
            {'text': '', 'font_size': Pt(14), 'font_color': WHITE, 'line_spacing': 1.0},
            {'text': 'Ne. Nikada.', 'font_size': Pt(28), 'font_color': WHITE, 'bold': True, 'line_spacing': 1.3, 'space_before': Pt(12)},
            {'text': '', 'font_size': Pt(8), 'font_color': WHITE, 'line_spacing': 1.0},
            {'text': 'Sve što izgovorite se obradi i odmah briše. Nijedan audio zapis ne ostaje pohranjen. Nikome.', 'font_size': Pt(15), 'font_color': ZINC_400, 'line_spacing': 1.5},
        ]
    )

    add_page_number(slide, 7, TOTAL_SLIDES)


def build_slide_08(prs):
    """SLIDE 08 — PLATFORMA IZBLIZA (Light)
    Fix: H1 44pt expanded box ("minuta" was clipped), phase descriptions from
    pitchdeck, removed [No Title] placeholder, mockup card labels fixed
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_logo(slide, is_dark=False)
    add_label(slide, "Kako izgledaju prvi koraci")

    # H1 — fix: 44pt, expanded height so "minuta" is fully visible
    add_multiline_text(
        slide, MARGIN, Inches(1.1), Inches(6.5), Inches(1.6),
        line_spacing=H1_LINE_HEIGHT,
        lines=[
            {'text': 'Spremni za rad', 'font_size': H1_SIZE_LIGHT, 'font_color': NEAR_BLACK, 'bold': True, 'line_spacing': H1_LINE_HEIGHT},
            {'text': 'za manje od 5 minuta.', 'font_size': H1_SIZE_LIGHT, 'font_color': ZINC_300, 'bold': True, 'line_spacing': H1_LINE_HEIGHT},
        ]
    )

    # Left — 5 phases (descriptions synced from canva_fix copy rewrites)
    phase_left = MARGIN
    phase_top = Inches(2.9)
    phase_width = Inches(5.5)
    phase_h = Inches(0.82)

    phases = [
        ("01", "Registracija",
         "Odete na aimed.cee-med.ba i registrujete se \u2014 kao što biste napravili email nalog. Traje 2 minute.",
         False),
        ("02", "Podaci ordinacije",
         "Unesete naziv vaše ordinacije, adresu i broj licence. Učitate sliku pečata i potpisa.",
         True),
        ("03", "Branding nalaza",
         "Dodajete logo \u2014 od tog trenutka svaki PDF koji izađe iz AiMED-a nosi vaš branding, automatski.",
         False),
        ("04", "Pacijentska baza",
         "Unesete pacijenta jednom. Svaki sljedeći nalaz za tu osobu vezuje se automatski uz njen karton.",
         False),
        ("05", "Odaberite mod",
         "Pritisnete 'Snimi' i počnete diktirati. Odaberete format i kliknete Export. Gotovo.",
         False),
    ]

    for i, (num, title, desc, is_accent) in enumerate(phases):
        py = phase_top + int(phase_h * i)

        # Circle
        circle_size = Inches(0.35)
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, phase_left, py + Inches(0.1), circle_size, circle_size
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = BLACK if not is_accent else WHITE
        circle.line.fill.background()

        add_text_box(slide, phase_left - Inches(0.01), py + Inches(0.07),
                     circle_size + Inches(0.02), circle_size + Inches(0.02),
                     text=num, font_size=Pt(11),
                     font_color=WHITE if not is_accent else BLACK,
                     bold=True, alignment=PP_ALIGN.CENTER, line_spacing=1.0)

        text_left = phase_left + Inches(0.55)

        if is_accent:
            add_rounded_rect(slide, text_left - Inches(0.1), py - Inches(0.02),
                             phase_width - Inches(0.3), phase_h + Inches(0.04),
                             fill_color=BLACK)

        add_multiline_text(
            slide, text_left, py + Inches(0.05),
            phase_width - Inches(0.8), phase_h - Inches(0.08),
            lines=[
                {'text': title, 'font_size': Pt(14), 'font_color': WHITE if is_accent else NEAR_BLACK, 'bold': True, 'line_spacing': 1.2},
                {'text': desc, 'font_size': Pt(12), 'font_color': ZINC_400 if is_accent else ZINC_600, 'line_spacing': 1.4},
            ]
        )

    # Right — Portal mockup cards (fix: proper labels, no "[No Title]")
    mock_left = Inches(7.0)
    mock_w = Inches(5.7)

    # Card 1 — Patient list (fix: "LISTA PACIJENATA" label)
    c1_top = Inches(2.6)
    c1_h = Inches(1.3)
    add_rounded_rect(slide, mock_left, c1_top, mock_w, c1_h,
                     fill_color=WHITE, border_color=ZINC_100)
    add_multiline_text(
        slide, mock_left + Inches(0.3), c1_top + Inches(0.15),
        mock_w - Inches(0.6), c1_h - Inches(0.3),
        lines=[
            {'text': 'LISTA PACIJENATA', 'font_size': Pt(10), 'font_color': ZINC_400, 'bold': True, 'line_spacing': 1.2, 'all_caps': True, 'spacing': LABEL_SPACING},
            {'text': 'Zdravko Čorić  |  ID: 29402  |  3 nalaza', 'font_size': Pt(13), 'font_color': NEAR_BLACK, 'line_spacing': 1.6, 'space_before': Pt(6)},
            {'text': 'Mina Kovač     |  ID: 18764  |  1 nalaz', 'font_size': Pt(13), 'font_color': NEAR_BLACK, 'line_spacing': 1.6},
        ]
    )

    # Card 2 — Recording (fix: "DIKTIRANJE U TOKU..." label)
    c2_top = c1_top + c1_h + Inches(0.15)
    c2_h = Inches(1.2)
    add_rounded_rect(slide, mock_left, c2_top, mock_w, c2_h, fill_color=BLACK)
    add_multiline_text(
        slide, mock_left + Inches(0.3), c2_top + Inches(0.2),
        mock_w - Inches(0.6), c2_h - Inches(0.4),
        alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        lines=[
            {'text': '\U0001F534  00:42 \u00b7 Diktiranje u toku...', 'font_size': Pt(18), 'font_color': WHITE, 'bold': True, 'line_spacing': 1.3},
            {'text': '\u2581\u2582\u2583\u2585\u2587\u2585\u2583\u2582\u2581\u2582\u2583\u2585\u2587\u2585\u2583\u2582\u2581\u2582\u2583\u2585\u2587\u2585\u2583\u2582\u2581', 'font_size': Pt(16), 'font_color': ZINC_500, 'line_spacing': 1.0, 'space_before': Pt(6)},
        ]
    )

    # Card 3 — Output (fix: "STRUKTURIRANI NALAZ" label)
    c3_top = c2_top + c2_h + Inches(0.15)
    c3_h = Inches(1.6)
    add_rounded_rect(slide, mock_left, c3_top, mock_w, c3_h,
                     fill_color=WHITE, border_color=ZINC_100)
    add_multiline_text(
        slide, mock_left + Inches(0.3), c3_top + Inches(0.15),
        mock_w - Inches(0.6), c3_h - Inches(0.3),
        lines=[
            {'text': 'STRUKTURIRANI NALAZ', 'font_size': Pt(10), 'font_color': ZINC_400, 'bold': True, 'line_spacing': 1.2, 'all_caps': True, 'spacing': LABEL_SPACING},
            {'text': 'Anamneza: Pacijent se žali na bolove u abdomenu...', 'font_size': Pt(12), 'font_color': NEAR_BLACK, 'line_spacing': 1.5, 'space_before': Pt(6)},
            {'text': 'Dijagnoza: K29.5 \u2014 Gastritis, nespecificiran', 'font_size': Pt(12), 'font_color': NEAR_BLACK, 'line_spacing': 1.5},
            {'text': 'Terapija: Sumamed 500mg 1x1, 3 dana', 'font_size': Pt(12), 'font_color': NEAR_BLACK, 'line_spacing': 1.5},
        ]
    )

    # Export buttons
    btn_y = c3_top + c3_h - Inches(0.45)
    for j, lbl in enumerate(["PDF", "Word", "Kopiraj"]):
        bx = mock_left + Inches(0.3) + int(Inches(1.3) * j)
        bw = Inches(1.1)
        add_rounded_rect(slide, bx, btn_y, bw, Inches(0.32),
                         fill_color=BLACK if j < 2 else ZINC_50,
                         border_color=ZINC_100 if j == 2 else None)
        add_text_box(slide, bx, btn_y + Inches(0.02), bw, Inches(0.28),
                     text=lbl, font_size=Pt(11),
                     font_color=WHITE if j < 2 else NEAR_BLACK,
                     bold=True, alignment=PP_ALIGN.CENTER, line_spacing=1.0)

    add_page_number(slide, 8, TOTAL_SLIDES)


def build_slide_faq(prs):
    """APPENDIX — FAQ SLAJD (Light)
    Fix: Q 14pt Bold black, A 13pt Regular zinc-600, fuller answers from pitchdeck,
    increased card height for readability (45+ year old doctors)
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_logo(slide, is_dark=False)
    add_label(slide, "Česta pitanja")

    add_text_box(slide, MARGIN, Inches(1.1), Inches(8), Inches(0.6),
                 text="Česta pitanja o AI diktiranju.",
                 font_size=Pt(36), font_color=NEAR_BLACK, bold=True, line_spacing=1.0)

    # 2 columns × 4 rows — fix: 14pt Q, 13pt A, 20px card gap
    col_width = Inches(5.8)
    card_h = Inches(1.2)        # Taller for readability
    card_gap = Inches(0.18)     # 20px gap
    col_gap = Inches(0.5)
    grid_top = Inches(2.1)

    # Fuller answers synced from aimed_pitchdeck.md
    questions = [
        # Left column
        [
            ("Kako početi sa AiMED-om?",
             "Kreirajte demo nalog, povežite mikrofon i započnite diktiranje. Sistem vas vodi kroz prvi proces strukturiranja nalaza."),
            ("Da li su podaci pacijenata sigurni?",
             "Apsolutno. End-to-end enkripcija, GDPR Član 9. Podaci se ne koriste za treniranje javnih modela. Vaši nalazi pripadaju isključivo vama."),
            ("Trebam li kupovati posebnu opremu?",
             "Ne. Mikrofon laptopa, slušalice od telefona ili pametni telefon. Vaš sadašnji uređaj je sve što vam treba."),
            ("Koliko traje obuka?",
             "Manje od 5 minuta. Interfejs je intuitivan čak i za ljekare koji nisu tehnički potkovani. Za klinike: demonstracija za tim."),
        ],
        # Right column
        [
            ("Koliko traje postavljanje sistema?",
             "Nula instalacije. Nakon profila, prvi diktat za manje od 60 sekundi. Radi odmah u Chrome, Edge, Safari."),
            ("Mogu li koristiti stare obrasce?",
             "Da. Učitajte vlastite šablone na koje ste navikli godinama. AiMED prilagođava tekst vašoj strukturi i brendiranju."),
            ("Mogu li segmentirati po pacijentima?",
             "Da. Nalazi su organizovani po pacijentu i datumu. Pretraga po imenu ili JMBG. Kronološki pregled svih nalaza."),
            ("Koliko često se AI modeli ažuriraju?",
             "Sedmično, na osnovu povratnih informacija ljekara iz regije. Preciznost se održava iznad 98%."),
        ]
    ]

    for col_idx, col_questions in enumerate(questions):
        col_left = MARGIN + int((col_width + col_gap) * col_idx)

        for row_idx, (q, a) in enumerate(col_questions):
            card_top = grid_top + int((card_h + card_gap) * row_idx)

            add_rounded_rect(slide, col_left, card_top, col_width, card_h,
                             fill_color=ZINC_50, border_color=ZINC_100)

            # Fix: Q 14pt Bold, A 13pt zinc-600, 16px spacing between
            add_multiline_text(
                slide, col_left + Inches(0.25), card_top + Inches(0.15),
                col_width - Inches(0.5), card_h - Inches(0.3),
                lines=[
                    {'text': q, 'font_size': Pt(14), 'font_color': NEAR_BLACK, 'bold': True, 'line_spacing': 1.3},
                    {'text': a, 'font_size': Pt(13), 'font_color': ZINC_600, 'line_spacing': 1.4, 'space_before': Pt(6)},
                ]
            )

    # Bottom strip
    strip_top = Inches(7.0)
    add_rounded_rect(slide, MARGIN, strip_top, SLIDE_W - MARGIN * 2, Inches(0.35),
                     fill_color=ZINC_50)
    add_text_box(slide, MARGIN, strip_top + Inches(0.04), SLIDE_W - MARGIN * 2, Inches(0.28),
                 text="Imate još pitanja? Kontaktirajte nas na info@cee-med.com",
                 font_size=Pt(13), font_color=ZINC_600, alignment=PP_ALIGN.CENTER, line_spacing=1.0)

    add_page_number(slide, 9, TOTAL_SLIDES)


def build_slide_09(prs):
    """SLIDE 09 — POZIV NA AKCIJU (Dark) — no font size change (dark slide)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BLACK)
    add_logo(slide, is_dark=True)

    # Label
    add_text_box(
        slide, Inches(0), Inches(1.2), SLIDE_W, Inches(0.3),
        text="SLJEDEĆI KORAK", font_size=LABEL_SIZE, font_color=ZINC_500,
        bold=True, alignment=PP_ALIGN.CENTER, all_caps=True, spacing=LABEL_SPACING, line_spacing=1.0
    )

    # H1
    add_multiline_text(
        slide, Inches(0), Inches(1.6), SLIDE_W, Inches(2.0),
        alignment=PP_ALIGN.CENTER, line_spacing=0.92,
        lines=[
            {'text': 'Zatražite pristup', 'font_size': Pt(72), 'font_color': WHITE, 'bold': True, 'line_spacing': 0.92},
            {'text': 'AiMED platformi.', 'font_size': Pt(72), 'font_color': WHITE, 'bold': True, 'line_spacing': 0.92},
        ]
    )

    # Subtitle
    add_multiline_text(
        slide, Inches(2.5), Inches(3.5), Inches(8.333), Inches(0.8),
        alignment=PP_ALIGN.CENTER,
        lines=[
            {'text': '14 dana besplatno. Bez kreditne kartice.', 'font_size': Pt(20), 'font_color': ZINC_400, 'line_spacing': 1.5},
            {'text': 'Personalizovana onboarding sesija za vašu ordinaciju.', 'font_size': Pt(20), 'font_color': ZINC_400, 'line_spacing': 1.5},
        ]
    )

    # 3 step cards
    card_top = Inches(4.5)
    card_w = Inches(3.8)
    card_h = Inches(1.1)
    card_gap = Inches(0.25)
    total_w = card_w * 3 + card_gap * 2
    start_x = (SLIDE_W - total_w) / 2

    steps = [
        ("01", "Zakažite besplatnu demonstraciju (30 min, online)"),
        ("02", "Probni period \u2014 14 dana besplatno, bez obaveza"),
        ("03", "Počnite štedjeti 1\u20132 sata dnevno od prvog dana"),
    ]

    for i, (num, text) in enumerate(steps):
        cx = int(start_x) + int((card_w + card_gap) * i)
        add_rounded_rect(slide, cx, card_top, card_w, card_h,
                         fill_color=DARK_SURFACE, border_color=RGBColor(0x27, 0x27, 0x2A))
        add_multiline_text(
            slide, cx + Inches(0.25), card_top + Inches(0.15),
            card_w - Inches(0.5), card_h - Inches(0.3),
            lines=[
                {'text': num, 'font_size': Pt(28), 'font_color': ZINC_700, 'bold': True, 'line_spacing': 1.1},
                {'text': text, 'font_size': Pt(13), 'font_color': ZINC_400, 'line_spacing': 1.4},
            ]
        )

    # Contact
    add_multiline_text(
        slide, Inches(0), Inches(5.9), SLIDE_W, Inches(1.0),
        alignment=PP_ALIGN.CENTER,
        lines=[
            {'text': '\u2709  info@cee-med.com', 'font_size': Pt(16), 'font_color': WHITE, 'line_spacing': 1.6},
            {'text': '\U0001F310  aimed.cee-med.ba', 'font_size': Pt(16), 'font_color': WHITE, 'line_spacing': 1.6},
            {'text': '\U0001F4CD  Sarajevo, Bosna i Hercegovina', 'font_size': Pt(16), 'font_color': WHITE, 'line_spacing': 1.6},
        ]
    )

    # Enterprise note
    add_text_box(
        slide, Inches(2), Inches(6.75), Inches(9.333), Inches(0.25),
        text="Enterprise integracije i grupne licence: info@cee-med.com",
        font_size=Pt(12), font_color=ZINC_500, alignment=PP_ALIGN.CENTER, line_spacing=1.0
    )

    # Footer
    add_text_box(
        slide, Inches(0), SLIDE_H - Inches(0.4), SLIDE_W, Inches(0.3),
        text="\u00a9 2026 AIMED | AI Medical Dictation \u00b7 CEE-MED d.o.o. \u00b7 Razvijeno za ljekare na Balkanu.",
        font_size=Pt(11), font_color=ZINC_700, alignment=PP_ALIGN.CENTER, line_spacing=1.0
    )

    add_page_number(slide, 10, TOTAL_SLIDES, is_dark=True)


# ═══════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Building AiMED Presentation v2 (synced with pitchdeck + fixes)...")

    builders = [
        ("01", "Naslovni slajd", build_slide_01),
        ("02", "Problem", build_slide_02),
        ("03", "Kontekst BiH", build_slide_03),
        ("04", "Rješenje", build_slide_04),
        ("05", "Kako radi", build_slide_05),
        ("06", "Funkcije", build_slide_06),
        ("07", "Sigurnost", build_slide_07),
        ("08", "Platforma izbliza", build_slide_08),
        ("09", "FAQ Appendix", build_slide_faq),
        ("10", "Poziv na akciju", build_slide_09),
    ]

    for i, (num, name, builder) in enumerate(builders, 1):
        print(f"  [{i:2d}/10] Slide {num} — {name}")
        builder(prs)

    output_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(output_dir, "AiMED_Prezentacija_2026.pptx")
    prs.save(output_path)
    print(f"\nPresentation saved to: {output_path}")
    print("Done!")


if __name__ == "__main__":
    main()
